import os
import time
import logging
from typing import Optional

logger = logging.getLogger(__name__)


def extract_text_from_file(file_path: str) -> Optional[str]:
    """
    Extract text from common document types for AI analysis.
    Supported: .pdf, .pptx, .docx, .txt, .md
    """
    if not os.path.exists(file_path):
        logger.error("File not found: %s", file_path)
        return None

    ext = os.path.splitext(file_path)[1].lower()
    start_time = time.time()

    try:
        if ext == '.pdf':
            logger.info("Extracting text from PDF: %s", file_path)
            text = _extract_from_pdf(file_path)
        elif ext == '.pptx':
            logger.info("Extracting text from PPTX: %s", file_path)
            text = _extract_from_pptx(file_path)
        elif ext == '.docx':
            logger.info("Extracting text from DOCX: %s", file_path)
            text = _extract_from_docx(file_path)
        elif ext in ['.txt', '.md']:
            logger.info("Reading text file: %s", file_path)
            text = _extract_from_text(file_path)
        else:
            logger.warning("Unsupported file type for extraction: %s", ext)
            return None

        logger.info("Text extraction complete in %.2fs (%d chars)", time.time() - start_time, len(text or ""))
        return text
    except Exception as exc:
        logger.error("Text extraction failed for %s: %s", file_path, exc)
        return None


def _extract_from_pdf(file_path: str) -> str:
    import pdfplumber

    text_content = []
    with pdfplumber.open(file_path) as pdf:
        for page_index, page in enumerate(pdf.pages, start=1):
            page_text = page.extract_text()

            if not page_text:
                # Some PDFs have extractable text that needs layout-aware parsing.
                page_text = page.extract_text(layout=True)

            if not page_text:
                # Fallback: rebuild simple text from extracted word tokens.
                words = page.extract_words(use_text_flow=True)
                if words:
                    page_text = " ".join(word.get("text", "") for word in words).strip()

            if page_text and page_text.strip():
                text_content.append(page_text)
            else:
                logger.debug("No text extracted from PDF page %d", page_index)

    full_text = "\n\n".join(text_content).strip()
    if not full_text:
        logger.warning(
            "PDF appears to have no machine-readable text (possibly scanned/image-only): %s",
            file_path,
        )
    return full_text


def _extract_from_pptx(file_path: str) -> str:
    import pptx

    text_content = []
    presentation = pptx.Presentation(file_path)
    for index, slide in enumerate(presentation.slides):
        text_content.append(f"--- Slide {index + 1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    text_content.append(" | ".join([cell.text_frame.text for cell in row.cells]))

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes:
                text_content.append(f"Notes: {notes}")

    return "\n".join(text_content)


def _extract_from_docx(file_path: str) -> str:
    import docx

    doc = docx.Document(file_path)
    return "\n".join([paragraph.text for paragraph in doc.paragraphs])


def _extract_from_text(file_path: str) -> str:
    with open(file_path, 'r', encoding='utf-8', errors='replace') as file_obj:
        return file_obj.read()
